"""
chunking.py — PDF text extraction and chunking engine.

Extracts text from large/messy PDFs via PyMuPDF (fitz), then splits
into configurable overlapping chunks suitable for LLM context windows.

All processing is in-memory — no intermediate files written.

Key behaviours:
  - Text cleaning happens PER PAGE before concatenation, so page-local
    artefacts (headers/footers) are stripped while page boundaries are
    still known.
  - Chunk boundaries are *soft-aligned*: the engine tries to break at
    paragraph (\n\n), newline, sentence-end (. ! ?), or whitespace —
    in that preference order — within a configurable search window.
  - Every chunk carries a deterministic ``chunk_id`` for traceability.
  - Chunks exceeding 1.5× the target size emit a warning.
"""

from __future__ import annotations

import hashlib
import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import fitz  # PyMuPDF

try:
    import docx as _docx          # python-docx
    _DOCX_AVAILABLE = True
except ImportError:
    _DOCX_AVAILABLE = False

try:
    import pptx as _pptx          # python-pptx
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

_OVERSIZED_FACTOR = 1.5  # warn if chunk exceeds this × chunk_size

# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class Chunk:
    """A single text chunk with provenance metadata."""

    text: str
    chunk_id: str              # deterministic hash-based id
    index: int                 # 0-based chunk sequence number
    page_start: int            # first PDF page contributing to this chunk
    page_end: int              # last  PDF page contributing to this chunk
    char_count: int = field(init=False)

    def __post_init__(self) -> None:
        object.__setattr__(self, "char_count", len(self.text))

    def to_dict(self) -> dict:
        return {
            "text": self.text,
            "chunk_id": self.chunk_id,
            "index": self.index,
            "page_start": self.page_start,
            "page_end": self.page_end,
            "char_count": self.char_count,
        }


@dataclass
class ChunkingConfig:
    """Tunable parameters for the chunker."""

    chunk_size: int = 4000          # target chars per chunk
    chunk_overlap: int = 400        # overlap chars between consecutive chunks
    min_chunk_size: int = 200       # merge trailing runt below this
    boundary_search_window: int = 200  # chars to look back for soft break
    strip_headers_footers: bool = True
    normalize_whitespace: bool = True

    def __post_init__(self) -> None:
        if self.chunk_overlap >= self.chunk_size:
            raise ValueError("chunk_overlap must be < chunk_size")
        if self.min_chunk_size < 0:
            raise ValueError("min_chunk_size must be >= 0")
        if self.boundary_search_window < 0:
            raise ValueError("boundary_search_window must be >= 0")


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class PageText:
    """Text from a single PDF page (after cleaning)."""

    page_num: int   # 0-based
    text: str


def extract_pages(
    pdf_path: str | Path,
    config: ChunkingConfig | None = None,
) -> list[PageText]:
    """Extract and clean text from every page of a PDF.

    Cleaning is applied per-page so page-local artefacts (headers,
    footers, standalone page numbers) are stripped while page context
    is still available.
    """
    if config is None:
        config = ChunkingConfig()
    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        raise FileNotFoundError(f"PDF not found: {pdf_path}")

    pages: list[PageText] = []
    doc = fitz.open(str(pdf_path))
    total_pages = len(doc)
    try:
        for page_num in range(total_pages):
            raw = doc[page_num].get_text("text")
            cleaned = _clean_page(raw, config)
            if cleaned:
                pages.append(PageText(page_num=page_num, text=cleaned))
            else:
                logger.debug("Page %d: empty after cleaning, skipped", page_num)
    finally:
        doc.close()

    logger.info(
        "Extracted text from %d / %d pages in %s",
        len(pages), total_pages, pdf_path.name,
    )
    return pages


# ---------------------------------------------------------------------------
# Per-page text cleaning
# ---------------------------------------------------------------------------

_MULTI_SPACE = re.compile(r"[ \t]+")
_MULTI_NEWLINE = re.compile(r"\n{3,}")
_HEADER_FOOTER = re.compile(
    r"^\s*(?:page\s+)?\d+(?:\s+of\s+\d+)?\s*$",
    re.IGNORECASE | re.MULTILINE,
)


def _clean_page(raw: str, config: ChunkingConfig) -> str:
    """Clean a single page's text in-memory."""
    text = raw
    if config.strip_headers_footers:
        text = _HEADER_FOOTER.sub("", text)
    if config.normalize_whitespace:
        text = _MULTI_SPACE.sub(" ", text)
        text = _MULTI_NEWLINE.sub("\n\n", text)
    return text.strip()


# ---------------------------------------------------------------------------
# Chunk ID generation
# ---------------------------------------------------------------------------


def _make_chunk_id(source_name: str, index: int, text: str) -> str:
    """Deterministic chunk ID: short hash of source + index + content prefix."""
    payload = f"{source_name}|{index}|{text[:256]}"
    digest = hashlib.sha256(payload.encode("utf-8")).hexdigest()[:12]
    return f"chunk-{index:04d}-{digest}"


# ---------------------------------------------------------------------------
# Soft boundary alignment
# ---------------------------------------------------------------------------

# Priority-ordered break patterns (best → worst)
_BREAK_PATTERNS: list[tuple[str, re.Pattern[str]]] = [
    ("paragraph", re.compile(r"\n\n")),
    ("newline",   re.compile(r"\n")),
    ("sentence",  re.compile(r"[.!?]\s")),
    ("space",     re.compile(r"\s")),
]


def _soft_break(text: str, hard_end: int, window: int) -> int:
    """Find the best break point near ``hard_end``.

    Searches backwards from ``hard_end`` within ``window`` chars for the
    highest-priority natural boundary.  Returns the offset (exclusive)
    at which to cut.  Falls back to ``hard_end`` if nothing found.
    """
    if hard_end >= len(text):
        return hard_end  # at document end, no adjustment needed

    search_start = max(0, hard_end - window)
    region = text[search_start:hard_end]

    for _name, pattern in _BREAK_PATTERNS:
        # Find the LAST match in the search region
        match = None
        for m in pattern.finditer(region):
            match = m
        if match is not None:
            # Cut after the matched break
            return search_start + match.end()

    return hard_end  # no natural break found


# ---------------------------------------------------------------------------
# Page-map helpers
# ---------------------------------------------------------------------------


def _build_page_map(pages: list[PageText]) -> tuple[str, list[tuple[int, int]]]:
    """Concatenate already-cleaned page texts and build offset → page_num map."""
    parts: list[str] = []
    spans: list[tuple[int, int]] = []
    offset = 0
    for pt in pages:
        spans.append((offset, pt.page_num))
        parts.append(pt.text)
        offset += len(pt.text) + 2  # +2 for "\n\n" separator
    return "\n\n".join(parts), spans


def _page_at_offset(offset: int, spans: list[tuple[int, int]]) -> int:
    """Binary-search page_spans to find which page owns an offset."""
    lo, hi = 0, len(spans) - 1
    while lo < hi:
        mid = (lo + hi + 1) // 2
        if spans[mid][0] <= offset:
            lo = mid
        else:
            hi = mid - 1
    return spans[lo][1]


# ---------------------------------------------------------------------------
# Chunking engine
# ---------------------------------------------------------------------------


def chunk_text(
    pages: list[PageText],
    config: ChunkingConfig | None = None,
    source_name: str = "unknown",
) -> list[Chunk]:
    """Split pre-cleaned pages into overlapping, soft-boundary-aligned chunks.

    Algorithm:
      1. Concatenate all (already cleaned) page texts.
      2. Walk forward by (chunk_size - overlap).
      3. At each raw boundary, search backwards within
         ``boundary_search_window`` for the best natural break.
      4. Annotate each chunk with source page range and chunk_id.
      5. Warn if any chunk exceeds 1.5× target size.

    Returns list of Chunk objects, fully in-memory.
    """
    if config is None:
        config = ChunkingConfig()

    if not pages:
        return []

    full_text, page_spans = _build_page_map(pages)
    if not full_text.strip():
        return []

    step = config.chunk_size - config.chunk_overlap
    warn_threshold = int(config.chunk_size * _OVERSIZED_FACTOR)
    chunks: list[Chunk] = []
    idx = 0
    pos = 0

    while pos < len(full_text):
        raw_end = min(pos + config.chunk_size, len(full_text))
        end = _soft_break(full_text, raw_end, config.boundary_search_window)
        slice_text = full_text[pos:end].strip()

        if not slice_text:
            break

        # Trailing runt — merge into previous chunk
        is_final = end >= len(full_text)
        if is_final and len(slice_text) < config.min_chunk_size and chunks:
            last = chunks[-1]
            merged_text = last.text + "\n" + slice_text
            chunks[-1] = Chunk(
                text=merged_text,
                chunk_id=last.chunk_id,
                index=last.index,
                page_start=last.page_start,
                page_end=_page_at_offset(
                    min(end - 1, len(full_text) - 1), page_spans
                ),
            )
            if len(merged_text) > warn_threshold:
                logger.warning(
                    "Chunk %d oversized after runt merge: %d chars "
                    "(threshold %d)",
                    last.index, len(merged_text), warn_threshold,
                )
            break

        page_start = _page_at_offset(pos, page_spans)
        page_end = _page_at_offset(
            min(end - 1, len(full_text) - 1), page_spans
        )
        chunk_id = _make_chunk_id(source_name, idx, slice_text)

        if len(slice_text) > warn_threshold:
            logger.warning(
                "Chunk %d oversized: %d chars (threshold %d)",
                idx, len(slice_text), warn_threshold,
            )

        chunks.append(Chunk(
            text=slice_text,
            chunk_id=chunk_id,
            index=idx,
            page_start=page_start,
            page_end=page_end,
        ))
        idx += 1

        # Advance by step, but account for soft-break adjustment:
        # next chunk starts at (end - overlap) so overlap region is preserved.
        pos = max(pos + step, end - config.chunk_overlap)

    logger.info(
        "Chunked %d chars into %d chunks (size=%d, overlap=%d) [%s]",
        len(full_text), len(chunks), config.chunk_size,
        config.chunk_overlap, source_name,
    )
    return chunks


# ---------------------------------------------------------------------------
# Convenience: single-call PDF → chunks
# ---------------------------------------------------------------------------


def chunk_pdf(
    pdf_path: str | Path,
    config: ChunkingConfig | None = None,
) -> list[Chunk]:
    """End-to-end: PDF file → list of Chunk objects, all in-memory."""
    pdf_path = Path(pdf_path)
    pages = extract_pages(pdf_path, config)
    return chunk_text(pages, config, source_name=pdf_path.stem)


# ---------------------------------------------------------------------------
# Convenience: single-call plain-text file → chunks
# ---------------------------------------------------------------------------


def chunk_text_file(
    file_path: str | Path,
    config: ChunkingConfig | None = None,
) -> list[Chunk]:
    """End-to-end: plain-text file → list of Chunk objects, all in-memory.

    The file is treated as a single virtual "page 0" so all downstream
    pipeline steps work unchanged.  Encoding is detected via UTF-8 with
    a latin-1 fallback.
    """
    if config is None:
        config = ChunkingConfig()
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    try:
        raw = file_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        raw = file_path.read_text(encoding="latin-1")

    cleaned = _clean_page(raw, config)
    if not cleaned:
        logger.warning("Text file %s is empty after cleaning", file_path.name)
        return []

    pages = [PageText(page_num=0, text=cleaned)]
    logger.info("Loaded text file %s (%d chars)", file_path.name, len(cleaned))
    return chunk_text(pages, config, source_name=file_path.stem)


# ---------------------------------------------------------------------------
# Word (.docx) extractor
# ---------------------------------------------------------------------------

def chunk_docx(
    file_path: str | Path,
    config: ChunkingConfig | None = None,
) -> list[Chunk]:
    """Extract text from a .docx file and chunk it."""
    if not _DOCX_AVAILABLE:
        raise ImportError("python-docx is required for .docx files: pip install python-docx")
    file_path = Path(file_path)
    doc = _docx.Document(str(file_path))
    # Each paragraph becomes a line; tables are also extracted
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_text = "  |  ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    raw = "\n\n".join(parts)
    if config is None:
        config = ChunkingConfig()
    cleaned = _clean_page(raw, config)
    if not cleaned:
        logger.warning("Docx file %s is empty after cleaning", file_path.name)
        return []
    pages = [PageText(page_num=0, text=cleaned)]
    logger.info("Loaded docx file %s (%d chars)", file_path.name, len(cleaned))
    return chunk_text(pages, config, source_name=file_path.stem)


# ---------------------------------------------------------------------------
# PowerPoint (.pptx) extractor
# ---------------------------------------------------------------------------

def chunk_pptx(
    file_path: str | Path,
    config: ChunkingConfig | None = None,
) -> list[Chunk]:
    """Extract text from a .pptx file and chunk it."""
    if not _PPTX_AVAILABLE:
        raise ImportError("python-pptx is required for .pptx files: pip install python-pptx")
    file_path = Path(file_path)
    prs = _pptx.Presentation(str(file_path))
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide_parts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_parts))
    raw = "\n\n".join(parts)
    if config is None:
        config = ChunkingConfig()
    cleaned = _clean_page(raw, config)
    if not cleaned:
        logger.warning("Pptx file %s is empty after cleaning", file_path.name)
        return []
    pages = [PageText(page_num=0, text=cleaned)]
    logger.info("Loaded pptx file %s (%d chars)", file_path.name, len(cleaned))
    return chunk_text(pages, config, source_name=file_path.stem)


# ---------------------------------------------------------------------------
# Convenience: auto-detect file type and chunk accordingly
# ---------------------------------------------------------------------------

_TEXT_SUFFIXES = frozenset({".txt", ".md", ".rst", ".csv", ".log"})


def chunk_file(
    file_path: str | Path,
    config: ChunkingConfig | None = None,
) -> list[Chunk]:
    """Route to the correct extractor based on file extension.

    Supported formats:
      .txt / .md / .rst / .csv / .log  — plain text
      .pdf                              — PDF (PyMuPDF)
      .docx                             — Word (python-docx)
      .pptx                             — PowerPoint (python-pptx)
    """
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()
    if suffix in _TEXT_SUFFIXES:
        return chunk_text_file(file_path, config)
    if suffix == ".docx":
        return chunk_docx(file_path, config)
    if suffix == ".pptx":
        return chunk_pptx(file_path, config)
    return chunk_pdf(file_path, config)
